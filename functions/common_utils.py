import re
import json
import base64
import io
from concurrent.futures import ThreadPoolExecutor
from PIL import Image
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import config


def clean_page_text_list(page_text_list: list[dict]) -> list[dict]:
    """
    删除每页文本中的 Markdown 图片：
    ![图片名.jpg](图片地址)
    """
    cleaned_list = []

    for item in page_text_list:
        text = item.get("text", "")

        # 删除 Markdown 图片格式，例如：
        # ![09e3-3.jpg](http://example.com/xxx.jpg)
        cleaned_text = re.sub(r"!\[[^\]]*\]\([^)]*\)", "", text)

        # 清理图片删除后多余的空行和首尾空格
        cleaned_text = re.sub(r"\n{3,}", "\n\n", cleaned_text).strip()

        cleaned_list.append({
            "page_num": item.get("page_num"),
            "text": cleaned_text,
        })

    return cleaned_list


def encode_image(image_path: str, high_resolution: bool = False) -> str:
    """图片转base64（默认压缩到1024x1024，high_resolution保留原图）"""
    with Image.open(image_path) as img:
        if not high_resolution:
            img = img.resize((1024, 1024), Image.LANCZOS)
        buffer = io.BytesIO()
        img.convert("RGB").save(buffer, format="JPEG")
        return base64.b64encode(buffer.getvalue()).decode("ascii")


def extract_ppt_images(ppt_path: str, download_dir: str = config.DOWNLOAD_DIR) -> list[dict]:
    """
    提取 PPT 每页中的图片。

    返回格式：
    [
        {1: ["/xxx/page_1_image_1.png"]},
        {2: ["/xxx/page_2_image_1.jpg", "/xxx/page_2_image_2.png"]}
    ]
    """
    file_name = Path(ppt_path).stem
    output_dir = Path(download_dir) / file_name
    output_dir.mkdir(parents=True, exist_ok=True)

    prs = Presentation(ppt_path)
    result = []

    for page_num, slide in enumerate(prs.slides, 1):
        image_paths = []

        for index, shape in enumerate(
            (s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE),
            1
        ):
            image = shape.image
            image_path = output_dir / f"page_{page_num}_image_{index}.{image.ext}"
            image_path.write_bytes(image.blob)
            image_paths.append(str(image_path))

        if image_paths:
            result.append({
                "page_num": page_num,
                "image_paths": image_paths
            })

    return result


def collect_multimedia_ocr(ppt_path: str) -> dict:
    """并行识别PPT中的图片与视频，返回按页分组的结果。"""
    # 延迟导入，避免循环依赖（prompt_message -> call_qwen -> common_utils）
    from prompt_message import ppt_image_describer, ppt_video_describer
    from functions.concurrent_util import execute_parallel_with_fallback
    from functions.video_extractor import process_ppt_videos

    def parse_result(raw: str) -> dict:
        try:
            data = json.loads(raw)
            return data if isinstance(data, dict) else {}
        except (json.JSONDecodeError, TypeError):
            return {}

    def ocr_image(image_path: str) -> dict:
        data = parse_result(ppt_image_describer(image_path))
        return {
            "raw_text": data.get("raw_text", ""),
            "summary": data.get("summary", ""),
        }

    def video_summary(summary: str, duration_seconds: float | None) -> dict:
        if duration_seconds is None:
            return {"summary": summary}
        total_seconds = max(1, round(duration_seconds))
        minutes, seconds = divmod(total_seconds, 60)
        duration = f"{minutes}分{seconds}秒" if minutes else f"{seconds}秒"
        return {"summary": f"视频时长约{duration}。{summary}"}

    def ocr_video(frame_paths: list[str], duration_seconds: float | None) -> dict:
        data = parse_result(ppt_video_describer(frame_paths))
        return video_summary(data.get("summary", ""), duration_seconds)

    def process_images() -> dict:
        tasks, task_pages = [], []
        for page in extract_ppt_images(ppt_path):
            for image_path in page["image_paths"]:
                tasks.append((ocr_image, (image_path,), {},
                              Path(image_path).name,
                              {"raw_text": "", "summary": "该图片暂时无法读取"}))
                task_pages.append(page["page_num"])
        results = execute_parallel_with_fallback(tasks, max_workers=3, timeout_seconds=180)
        grouped = {}
        for page_num, result in zip(task_pages, results):
            grouped.setdefault(page_num, []).append(result)
        return grouped

    def process_videos() -> dict:
        video_dir = Path(config.DOWNLOAD_DIR) / Path(ppt_path).stem
        tasks, task_pages = [], []
        for item in process_ppt_videos(ppt_path, output_dir=str(video_dir)):
            page_num = item["slide"]
            if not page_num:
                print(f"⚠️  视频无法定位页码，跳过: {Path(item['video']).name}")
                continue
            frame_paths = [path for path in item["frames"].values() if path]
            if not frame_paths:
                continue
            duration_seconds = item.get("duration_seconds")
            tasks.append((ocr_video, (frame_paths, duration_seconds), {},
                          Path(item["video"]).name,
                          video_summary("该视频暂时无法读取", duration_seconds)))
            task_pages.append(page_num)
        results = execute_parallel_with_fallback(tasks, max_workers=3, timeout_seconds=300)
        grouped = {}
        for page_num, result in zip(task_pages, results):
            grouped.setdefault(page_num, []).append(result)
        return grouped

    with ThreadPoolExecutor(max_workers=2) as executor:
        images_future = executor.submit(process_images)
        videos_future = executor.submit(process_videos)
        return {
            "images_by_page": images_future.result(),
            "videos_by_page": videos_future.result(),
        }


def enrich_pages_with_ocr(
    ppt_path: str,
    page_text_list: list[dict],
    multimedia_result: dict | None = None,
) -> list[dict]:
    """清理页面正文并合并结构化图片、视频识别结果。"""
    if multimedia_result is None:
        multimedia_result = collect_multimedia_ocr(ppt_path)
    images_by_page = multimedia_result.get("images_by_page", {})
    videos_by_page = multimedia_result.get("videos_by_page", {})

    enriched = clean_page_text_list(page_text_list)
    for item in enriched:
        page_num = item["page_num"]
        item["images"] = images_by_page.get(page_num, [])
        item["videos"] = videos_by_page.get(page_num, [])

    return enriched

if __name__ == "__main__":
    print(extract_ppt_images(r"C:\Users\shen.xin\Downloads\AI&财务\AI智算中心&财务P1工作.pptx"))
